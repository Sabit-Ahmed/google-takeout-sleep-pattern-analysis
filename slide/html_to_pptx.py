from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_summary_slide():
    # Create presentation
    prs = Presentation()

    # Use 16:9 slide layout
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_text = title.text_frame
    title_text.text = "Digital Activity Analysis"
    title_para = title_text.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(14), Inches(0.5))
    subtitle_text = subtitle.text_frame
    subtitle_text.text = "Data Collection Period: Jan 01 - Feb 13, 2025 (1,859 entries)"
    subtitle_para = subtitle_text.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(20)

    # Left column - Statistics
    stats = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
    stats_frame = stats.text_frame

    p = stats_frame.add_paragraph()
    p.text = "Key Statistics"
    p.font.size = Pt(28)
    p.font.bold = True

    # Daily Patterns
    p = stats_frame.add_paragraph()
    p.text = "Daily Patterns"
    p.font.size = Pt(20)
    p.font.bold = True

    for item in ["Peak Hours: 23:00, 16:00, 17:00",
                 "Quietest: 02:00-07:59",
                 "Average Daily Activities: 42.2"]:
        p = stats_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)

    # Sleep Patterns
    p = stats_frame.add_paragraph()
    p.text = "\nSleep Patterns"
    p.font.size = Pt(20)
    p.font.bold = True

    for item in ["Weekday Mean: 8.50h",
                 "Weekend Mean: 9.61h",
                 "No significant difference (p=0.0677)"]:
        p = stats_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)

    # Middle column - placeholder for weekly pattern
    weekly = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(5))
    weekly_frame = weekly.text_frame
    p = weekly_frame.add_paragraph()
    p.text = "Weekly Activity Pattern"
    p.font.size = Pt(28)
    p.font.bold = True

    p = weekly_frame.add_paragraph()
    p.text = "[Insert Figure 4 here]"
    p.font.size = Pt(16)
    p.font.italic = True

    # Right column - Topic Analysis
    topic = slide.shapes.add_textbox(Inches(9.5), Inches(2), Inches(4), Inches(5))
    topic_frame = topic.text_frame
    p = topic_frame.add_paragraph()
    p.text = "Topic Analysis"
    p.font.size = Pt(28)
    p.font.bold = True

    p = topic_frame.add_paragraph()
    p.text = "[Insert Figure 6 here]"
    p.font.size = Pt(16)
    p.font.italic = True

    p = topic_frame.add_paragraph()
    p.text = "\n[Insert Figure 7 here]"
    p.font.size = Pt(16)
    p.font.italic = True

    # Save the presentation
    prs.save('project1_slide.pptx')


if __name__ == "__main__":
    create_summary_slide()